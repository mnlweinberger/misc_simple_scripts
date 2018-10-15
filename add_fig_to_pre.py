#!/usr/bin/python

import os, sys

from pptx import Presentation
from pptx.util import Inches

# Open a file
path = "/home/mor/Desktop/287_tRNA/"
dirs = os.listdir( path )
prs = Presentation()

# This would print all the files and directories
for pics in dirs:
	img_path = '/home/mor/Desktop/287_tRNA/'+pics
	print img_path
	
	blank_slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(blank_slide_layout)
	left = Inches(1)
	top = Inches(2)
	height = Inches(5)
	slide.shapes.add_picture(img_path, left, top, height=height)
	shapes = slide.shapes
	title_shape = shapes.title
	title_shape.text=pics


prs.save('/home/mor/Desktop/287_tRNA/test.pptx')
